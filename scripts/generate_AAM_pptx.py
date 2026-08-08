# -*- coding: utf-8 -*-
"""
generate_AAM_pptx.py

生成“先进空中移动（AAM）发展面临的安全挑战”1小时演讲 PPTX 的脚本。

说明：在本地运行此脚本会产生一个可打开的 PowerPoint 文件：docs/AAM-safety-1hr.pptx
依赖：python-pptx
  pip install python-pptx

用法：
  python scripts/generate_AAM_pptx.py

输出：
  docs/AAM-safety-1hr.pptx

若需要定制封面（演讲者/单位/日期）或样式，请编辑脚本顶部的 METADATA 字段。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os

# 可配置元数据（请按需修改）
METADATA = {
    "title": "先进空中移动（AAM）发展面临的安全挑战",
    "subtitle": "风险、缓解与实施路线（1 小时演讲）",
    "presenter": "XXX",
    "organization": "XXX",
    "date": "2026-08-08",
    "output_path": "docs/AAM-safety-1hr.pptx",
}

# 幻灯片内容：每项包含 title, bullets (list), notes (讲稿)
SLIDES = [
    {"title": "封面", "bullets": [METADATA['title'], METADATA['subtitle'], f"演讲者：{METADATA['presenter']}  单位：{METADATA['organization']}  日期：{METADATA['date']}"],
     "notes": f"大家好，我是{METADATA['presenter']}，来自{METADATA['organization']}。今天用一小时的时间，和大家一起探讨AAM在商业化与规模化过程中面对的主要安全挑战、缓解措施，以及可执行的路线图。"},
    {"title": "执行摘要", "bullets": ["AAM 的潜力", "主要安全挑战的总体类别", "演讲目标：识别风险、提出缓解、给出路线图"],
     "notes": "AAM通过 eVTOL 等新型飞行器承诺提供高效出行，但在车辆、自动化、空域、基础设施与社会接受度等方面引入新的风险。今天会识别这些风险并提出缓解策略。"},
    {"title": "议程", "bullets": ["背景", "风险领域", "缓解策略", "研究空白", "测试与路线图", "Q&A"],
     "notes": "接下来按照这个顺序展开，预计50分钟讲演+10分钟答疑。"},
    {"title": "空域整合：问题描述", "bullets": ["混合运行", "低空高密度", "城市航廊与短跳航段"],
     "notes": "AAM引入大量短航程、低空飞行器，空域从少量高空航路扩展到复杂的低空网络，带来混合运行冲突与更高的态势感知需求。"},
    {"title": "空域整合：流量管理对策", "bullets": ["分层 UTM 架构", "地理围栏与动态限值", "战略与战术解冲突"],
     "notes": "需要分层的流量管理：战略层做容量与航路规划，战术层做实时解冲突，并结合地理围栏与动态限制。"},
    {"title": "飞行器可靠性：挑战与失效模式", "bullets": ["电池/推进失效与热失控", "分布式推进的耦合失效", "结构疲劳与短航程循环"],
     "notes": "电气化与新架构带来新的失效模式，需要重构冗余与维护策略。"},
    {"title": "冗余与维护实践", "bullets": ["主动/被动冗余设计", "降级飞行包线与故障恢复", "基于使用的维护与检测制度"],
     "notes": "多重冗余、降级控制律与基于状态的维护计划可以降低单点失效风险。"},
    {"title": "自主感知：传感器挑战", "bullets": ["雨/雾/眩光下的退化", "城市背景下虚警/漏检", "多模态传感融合需求"],
     "notes": "感知系统在复杂环境下性能下降，需传感融合与感知退化检测。"},
    {"title": "自主决策与人机交互", "bullets": ["可验证与可解释的算法", "控制移交机制", "监督与远端操作员界面"],
     "notes": "自动化决策必须可验证并在人工交接时保持可预测行为。"},
    {"title": "人因：培训与操作", "bullets": ["远端/监督飞行员角色", "工作负荷与程序复杂性", "认证与培训体系"],
     "notes": "需要为新角色设计专门的培训与认证流程，以减少人为错误。"},
    {"title": "地面与维护人因", "bullets": ["换电/充电作业风险", "地面安全区与流程标准", "工具与接口的人因设计"],
     "notes": "地面操作同样高风险，应通过流程与工具降低事故概率。"},
    {"title": "基础设施：垂直起降点设计", "bullets": ["选址与进近/离场通道", "障碍物净空与应急通道", "消防/医疗接入"],
     "notes": "vertiport 需兼顾进近路径、障碍物及紧急救援接入。"},
    {"title": "电力与充电基础设施", "bullets": ["高吞吐充电与热管理", "电力质量与微电网设计", "储能与故障隔离"],
     "notes": "充电对电网与本地供电的需求需提前规划，包括储能和热管理措施。"},
    {"title": "天气与城市环境影响", "bullets": ["低空气象与风切变", "城市峡谷湍流", "短时高分辨率预报需求"],
     "notes": "城市微气候会显著影响低空性能与运营规则。"},
    {"title": "环境退化与寿命周期", "bullets": ["尘土/盐雾对传感器与推进的影响", "电池长期老化与循环特性", "疲劳与耐久性试验"],
     "notes": "长期耐久性研究需真实运营数据支撑的老化与疲劳试验。"},
    {"title": "网络安全：攻击面与威胁模型", "bullets": ["C2 与遥测链路风险", "OTA 与软件供应链", "入侵检测与响应"],
     "notes": "应从设计层面强化链路加密、认证与软件签名，并部署监测与响应体系。"},
    {"title": "数据完整性与 UTM 信任", "bullets": ["位置/意图伪造风险", "端到端认证與加密", "隐私保护与可溯源性"],
     "notes": "UTM 运行依赖可信数据交换，应采用认证与隐私保护机制。"},
    {"title": "认证与监管差距", "bullets": ["现有认证路径不匹配", "AI/ML 认证空白", "跨域监管协同难题"],
     "notes": "建议采用基于性能的认证方法并加强跨机构协作。"},
    {"title": "研究空白与优先方向", "bullets": ["AI/ML 认证方法学", "城市低空气象模拟", "可扩展且安全的 UTM", "电池老化长期研究", "人因学研究"],
     "notes": "列出优先研究主题并建议建立共享试验平台以加速验证。"},
    {"title": "风险缓解策略（总结）", "bullets": ["系统级安全工程与安全案例", "稳健冗余與渐进降级", "可认证自主系统", "网络安全工程", "以人为本的人因工程"],
     "notes": "将策略整合到安全案例中，通过持续测试验证有效性。"},
    {"title": "测试与验证计划（checklist）", "bullets": ["数字孪生与大规模仿真", "硬件在环 (HIL) 与传感器在环", "分阶段实飞验证", "场景库与边界测试"],
     "notes": "建立从仿真到实飞的分阶段验证体系，并构建场景库以覆盖边界情况。"},
    {"title": "实施路线图（3 年 / 5 年）", "bullets": ["概念验证 → 示范运营 → 扩规模部署", "监管对接与早期社区参与", "迭代验证与扩展"],
     "notes": "给出逐步实施里程碑，强调监管与社区的早期互动。"},
    {"title": "主要结论与行动建议", "bullets": ["建立安全案例模板", "启动传感退化与电池老化试验", "与监管共建验证方案"],
     "notes": "提出3-5条立即可执行的行动项，便于团队快速启动。"},
    {"title": "问答", "bullets": ["开放提问（10 分钟）"],
     "notes": "邀请听众提问，并准备2-3个备选问题以引导讨论。"},
]

# 简单的样式辅助函数

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    subtitle_tf = slide.placeholders[1]
    title_tf.text = title
    subtitle_tf.text = subtitle
    return slide


def add_bulleted_slide(prs, title, bullets, notes_text=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title_tf.text = title
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
    # 添加备注（讲稿）
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text
    return slide


def generate_pptx(output_path):
    prs = Presentation()
    # 首页使用 title slide
    first = SLIDES[0]
    add_title_slide(prs, first['title'], "\n".join(first['bullets']))
    prs.slides[0].notes_slide.notes_text_frame.text = first.get('notes','')

    for s in SLIDES[1:]:
        add_bulleted_slide(prs, s['title'], s['bullets'], s.get('notes',''))

    # 确保输出目录存在
    out_dir = os.path.dirname(output_path)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir)

    prs.save(output_path)
    print(f"已生成 PPTX：{output_path}")


if __name__ == '__main__':
    generate_pptx(METADATA['output_path'])
